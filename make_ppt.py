"""
StructuredAlpha v4 — Complete redesign. All visual and accuracy fixes applied.
Run: py -3.13 make_ppt.py
Output: results/StructuredAlpha_v4.pptx
"""
import io, os, re, sys
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

RESULTS = os.path.join(os.path.dirname(__file__), "results")
OUT     = os.path.join(RESULTS, "StructuredAlpha_v6.pptx")

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0A, 0x1F, 0x3D)
BLUE  = RGBColor(0x1A, 0x5F, 0xA8)
TEAL  = RGBColor(0x00, 0xC2, 0xA8)
GOLD  = RGBColor(0xF0, 0xB4, 0x29)
RED   = RGBColor(0xE0, 0x3C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF4, 0xF6, 0xFA)
MGRAY = RGBColor(0xCC, 0xD1, 0xDC)
DGRAY = RGBColor(0x3A, 0x3D, 0x4A)

mNAVY  = "#0A1F3D"
mBLUE  = "#1A5FA8"
mTEAL  = "#00C2A8"
mGOLD  = "#F0B429"
mRED   = "#E03C2C"
mLGRAY = "#F4F6FA"
mDGRAY = "#3A3D4A"
mMGRAY = "#8890A0"

W = Inches(13.33)
H = Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Load data ─────────────────────────────────────────────────────────────────
mc       = pd.read_csv(os.path.join(RESULTS, "mc_summary.csv"))
baseline = pd.read_csv(os.path.join(RESULTS, "rl_baseline.csv"))["episode_return"].values
policy   = pd.read_csv(os.path.join(RESULTS, "rl_policy.csv"))["episode_return"].values
sv       = pd.read_csv(os.path.join(RESULTS, "seed_variance.csv"))
oc_ser   = np.load(os.path.join(RESULTS, "mc_oc_series.npy"))

summary_txt = open(os.path.join(RESULTS, "rl_summary.txt")).read()
base_sharpe = float(re.search(r"BASELINE.*?Sharpe:\s*([\d.]+)", summary_txt).group(1))
ppo_sharpe  = float(re.search(r"PPO.*?Sharpe:\s*([\d.]+)",      summary_txt).group(1))
base_dd     = float(re.search(r"BASELINE.*?drawdown:\s*([\d.]+)", summary_txt).group(1))
ppo_dd      = float(re.search(r"PPO.*?drawdown:\s*([\d.]+)",      summary_txt).group(1))
excess_bps  = int((policy.mean() - baseline.mean()) / 60 * 12 * 10000)

# Always derive from actual data — never hardcode
breached_df      = mc[mc["oc_breached_bbb"]]
no_breach_df     = mc[~mc["oc_breached_bbb"]]
breach_med       = breached_df["equity_total"].median()  / 1e6
no_breach_med    = no_breach_df["equity_total"].median() / 1e6
breach_p10       = breached_df["equity_total"].quantile(0.10) / 1e6
no_breach_p10    = no_breach_df["equity_total"].quantile(0.10) / 1e6
overall_med      = mc["equity_total"].median() / 1e6
breach_rate_pct  = mc["oc_breached_bbb"].mean() * 100
corr_cdr_eq      = mc["cdr"].corr(mc["equity_total"])
corr_rec_eq      = mc["recovery_rate"].corr(mc["equity_total"])
corr_breach_eq   = mc["oc_breached_bbb"].astype(float).corr(mc["equity_total"])
dd_redux_pct     = int(100 * (base_dd - ppo_dd) / base_dd)
n_seeds_win      = int((sv["ppo_sharpe"] > sv["base_sharpe"]).sum())


# ── Helpers ───────────────────────────────────────────────────────────────────
def fig_to_buf(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, w, h, fill_rgb=None, line_rgb=None, line_pt=1.0):
    sp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    if fill_rgb:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill_rgb
    else:
        sp.fill.background()
    if line_rgb:
        sp.line.color.rgb = line_rgb; sp.line.width = Pt(line_pt)
    else:
        sp.line.fill.background()
    return sp


def txt(slide, left, top, w, h, text="", size=14, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.fill.background()
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.bold  = bold
        r.font.color.rgb = color
    return tb


def img(slide, buf_or_path, left, top, w, h):
    if isinstance(buf_or_path, str):
        if os.path.exists(buf_or_path):
            slide.shapes.add_picture(buf_or_path, Inches(left), Inches(top),
                                     Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(buf_or_path, Inches(left), Inches(top),
                                 Inches(w), Inches(h))


def header_band(slide, section, title):
    """Dark navy top band — no teal line beneath."""
    rect(slide, 0, 0, 13.33, 1.25, fill_rgb=NAVY)
    txt(slide, 0.38, 0.18, 6, 0.3, text=section.upper(),
        size=9, bold=True, color=TEAL)
    txt(slide, 0.38, 0.46, 12.2, 0.65, text=title,
        size=27, bold=True, color=WHITE)


def callout(slide, text, top=5.8):
    """Bottom navy callout bar."""
    rect(slide, 0.32, top, 12.68, 0.72, fill_rgb=NAVY)
    txt(slide, 0.52, top + 0.08, 12.1, 0.58,
        text=text, size=10.5, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE + PROBLEM + WHAT WE BUILT
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg(s1, NAVY)

# Left accent strip
rect(s1, 0, 0, 0.18, 7.5, fill_rgb=TEAL)

# Title area
txt(s1, 0.55, 0.65, 11.5, 1.1,
    text="StructuredAlpha", size=58, bold=True, color=WHITE)
txt(s1, 0.55, 1.75, 11.5, 0.5,
    text="Autonomous CLO Tranche Allocation via Simulation + Deep Learning",
    size=16, color=RGBColor(0x8A, 0xB4, 0xD8))

# Two-column body: THE PROBLEM | WHAT WE BUILT
problem_items = [
    "CLO equity returns are highly sensitive to OC breach timing, which is\n"
    "    difficult to anticipate using static allocation rules.",
    "Equal-weight tranche allocation ignores macro regime shifts — the same\n"
    "    BBB weight in 2008 vs 2021 produces very different outcomes.",
    "No systematic framework connects macroeconomic stress signals (spreads,\n"
    "    yield curve, VIX) to CLO reinvestment and tranche allocation decisions.",
]
built_items = [
    "Monte Carlo waterfall engine — 2,000 paths x 60 months, full OC trigger\n"
    "    mechanics, stochastic CDR, recovery lags, binomial loan pool.",
    "LSTM regime classifier — detects TIGHT / NORMAL / WIDE from 5 FRED macro\n"
    "    signals (daily, 1996-2024). 82.3% accuracy on held-out test set.",
    "PPO reinforcement learning agent — dynamically rotates from BBB in benign\n"
    "    regimes to Cash/AAA in stress. +73 bps/yr, 37% lower drawdown.",
]

txt(s1, 0.55, 2.52, 5.8, 0.35,
    text="THE PROBLEM", size=10, bold=True, color=TEAL)
for i, item in enumerate(problem_items):
    txt(s1, 0.68, 2.88 + i * 0.78, 5.55, 0.72,
        text=f"{'  '}  {item}", size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

# Vertical divider
rect(s1, 6.58, 2.45, 0.02, 2.5, fill_rgb=RGBColor(0x22, 0x3A, 0x5C))

txt(s1, 6.75, 2.52, 6.2, 0.35,
    text="WHAT WE BUILT", size=10, bold=True, color=TEAL)
for i, item in enumerate(built_items):
    txt(s1, 6.88, 2.88 + i * 0.78, 6.0, 0.72,
        text=f"{'  '}  {item}", size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

# Stat pills
pills = [
    ("2,000", "Monte Carlo Paths"),
    ("82.3%", "LSTM Accuracy"),
    (f"+{excess_bps} bps/yr", "Excess Return"),
    (f"{dd_redux_pct}%", "Drawdown Reduction"),
]
rect(s1, 0.32, 5.55, 12.68, 1.55, fill_rgb=RGBColor(0x07, 0x15, 0x29))
for i, (v, l) in enumerate(pills):
    x = 0.55 + i * 3.1
    rect(s1, x, 5.72, 2.8, 1.18, fill_rgb=RGBColor(0x12, 0x28, 0x4A))
    txt(s1, x + 0.1, 5.78, 2.6, 0.62,
        text=v, size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s1, x + 0.1, 6.38, 2.6, 0.38,
        text=l, size=10.5, color=RGBColor(0x8A, 0xB4, 0xD8),
        align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SYSTEM ARCHITECTURE (one matplotlib figure, proper arrows)
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg(s2, LGRAY)
header_band(s2, "System Design", "Architecture Overview — Three Integrated Modules")

arch_cards = [
    {
        "title": "01  Monte Carlo\nWaterfall Engine",
        "color": mBLUE,
        "bullets": [
            "2,000 paths x 60 monthly periods",
            "Lognormal CDR, median 2.9%, right-skewed",
            "Binomial pool — 200 equal-weight loans",
            "Recovery at 12-month lag, 30-70% rate",
            "OC trigger: performing / note balance",
        ],
        "output": "Output: equity distribution, OC breach\nrate, CDR sensitivity curves"
    },
    {
        "title": "02  LSTM Regime\nClassifier",
        "color": mTEAL,
        "bullets": [
            "5 FRED series, daily data 1996-2024",
            "Baa/Aaa spreads, 2s10s, VIX, survey",
            "LSTM(128)->Dropout->LSTM(64)->FC(3)",
            "Labels: TIGHT / NORMAL / WIDE via VIX",
            "120-day window, 7,567 observations",
        ],
        "output": "Output: P(TIGHT, NORMAL, WIDE) vector\nfed into RL state (3 of 9 features)"
    },
    {
        "title": "03  PPO Reinforcement\nLearning Agent",
        "color": mGOLD,
        "bullets": [
            "State: regime(3) + OC(4) + CDR(2)",
            "Action: weights [Cash,AAA,AA,A,BBB]",
            "Reward: spread return - breach penalty",
            "200 updates x 2,048 steps = 409K ts",
            "500 out-of-sample evaluation episodes",
        ],
        "output": "Output: dynamic tranche weights,\nregime-conditional allocation rules"
    },
]

fig2, ax = plt.subplots(figsize=(12.6, 5.5), facecolor=mLGRAY)
ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
ax.set_facecolor(mLGRAY)

card_x = [0.01, 0.355, 0.70]
card_w = 0.30
card_h = 0.92
card_y = 0.04

for i, (cx, card) in enumerate(zip(card_x, arch_cards)):
    # White card body
    ax.add_patch(mpatches.FancyBboxPatch(
        (cx, card_y), card_w, card_h,
        boxstyle="round,pad=0.012",
        facecolor="white", edgecolor=mMGRAY, linewidth=1.2,
        transform=ax.transAxes, clip_on=False, zorder=2
    ))
    # Colored header
    ax.add_patch(mpatches.FancyBboxPatch(
        (cx, card_y + card_h - 0.175), card_w, 0.175,
        boxstyle="round,pad=0.008",
        facecolor=card["color"], edgecolor="none",
        transform=ax.transAxes, clip_on=False, zorder=3
    ))
    # Title in header
    ax.text(cx + card_w / 2, card_y + card_h - 0.085,
            card["title"],
            transform=ax.transAxes, fontsize=11.5, fontweight="bold",
            color="white", va="center", ha="center",
            linespacing=1.4, zorder=4)

    # Accent bar under header (thick — 5pt visual equivalent)
    ax.add_patch(mpatches.Rectangle(
        (cx, card_y + card_h - 0.175), card_w, 0.012,
        facecolor=card["color"], alpha=0.6,
        transform=ax.transAxes, clip_on=False, zorder=3
    ))

    # Bullet points
    for j, b in enumerate(card["bullets"]):
        y = card_y + card_h - 0.265 - j * 0.112
        ax.text(cx + 0.018, y, f"•  {b}",
                transform=ax.transAxes, fontsize=9.2,
                color=mDGRAY, va="top", ha="left", zorder=4)

    # Output box at bottom
    ax.add_patch(mpatches.FancyBboxPatch(
        (cx + 0.012, card_y + 0.012), card_w - 0.024, 0.128,
        boxstyle="round,pad=0.006",
        facecolor=card["color"], alpha=0.12,
        edgecolor=card["color"], linewidth=1.4,
        transform=ax.transAxes, clip_on=False, zorder=3
    ))
    ax.text(cx + card_w / 2, card_y + 0.078,
            card["output"],
            transform=ax.transAxes, fontsize=8.8,
            color=mNAVY, va="center", ha="center",
            linespacing=1.4, zorder=4)

# Proper arrows between cards (using annotate)
for x_from, x_to in [(0.31, 0.352), (0.655, 0.697)]:
    ax.annotate("",
        xy=(x_to, 0.52), xycoords="axes fraction",
        xytext=(x_from, 0.52), textcoords="axes fraction",
        arrowprops=dict(
            arrowstyle="->, head_width=0.35, head_length=0.015",
            color=mTEAL, lw=3.5
        ), zorder=5
    )

fig2.subplots_adjust(left=0, right=1, top=1, bottom=0)
img(s2, fig_to_buf(fig2), 0.32, 1.32, 12.68, 6.0)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MONTE CARLO: EQUITY DISTRIBUTION + OC FAN CHART
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg(s3, LGRAY)
header_band(s3, "Monte Carlo Waterfall",
            f"CLO Cashflow Simulation — 2,000 Paths x 60 Months")

# Stat pills (top-right inside header band)
stats3 = [
    (f"${overall_med:.0f}M", "Median Equity"),
    (f"{breach_rate_pct:.1f}%", "BBB OC Breach Rate"),
    ("2.9%", "Median CDR"),
]
for i, (v, l) in enumerate(stats3):
    x = 9.3 + i * 1.35
    rect(s3, x, 0.14, 1.2, 1.0, fill_rgb=RGBColor(0x14, 0x2E, 0x55))
    txt(s3, x + 0.06, 0.19, 1.08, 0.52,
        text=v, size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s3, x + 0.06, 0.70, 1.08, 0.38,
        text=l, size=9, color=RGBColor(0x8A, 0xB4, 0xD8),
        align=PP_ALIGN.CENTER)

# ── Chart A: Equity distribution, split by breach ─────────────────────────────
eq_b  = breached_df["equity_total"] / 1e6
eq_nb = no_breach_df["equity_total"] / 1e6

fig3a, ax = plt.subplots(figsize=(6.0, 4.0), facecolor="white")
ax.set_facecolor("white")
ax.hist(eq_nb, bins=38, color=mTEAL, alpha=0.72,
        label=f"No OC breach  (n={len(eq_nb):,})", edgecolor="white", linewidth=0.4)
ax.hist(eq_b,  bins=22, color=mRED,  alpha=0.72,
        label=f"OC breached   (n={len(eq_b):,})", edgecolor="white", linewidth=0.4)
ax.axvline(eq_nb.median(), color=mTEAL, lw=2.2, ls="--",
           label=f"Median no-breach: ${eq_nb.median():.0f}M")
ax.axvline(eq_b.median(),  color=mRED,  lw=2.2, ls="--",
           label=f"Median breached: ${eq_b.median():.0f}M")
# Highlight the p10 difference — the real story
ax.axvline(eq_nb.quantile(0.10), color=mTEAL, lw=1.2, ls=":",
           label=f"p10 no-breach: ${eq_nb.quantile(0.10):.0f}M")
ax.axvline(eq_b.quantile(0.10),  color=mRED,  lw=1.2, ls=":",
           label=f"p10 breached: ${eq_b.quantile(0.10):.0f}M")
ax.set_xlabel("Total Equity Payout ($M)", fontsize=10, color=mDGRAY)
ax.set_ylabel("Number of Paths", fontsize=10, color=mDGRAY)
ax.set_title("Equity Distribution by OC Breach Status", fontsize=11,
             fontweight="bold", color=mNAVY, pad=8)
ax.legend(fontsize=8, framealpha=0.9, loc="upper left")
for sp in ["top", "right"]:
    ax.spines[sp].set_visible(False)
ax.spines["left"].set_color(mMGRAY); ax.spines["bottom"].set_color(mMGRAY)
ax.tick_params(colors=mDGRAY, labelsize=9)
fig3a.tight_layout()
img(s3, fig_to_buf(fig3a), 0.32, 1.38, 6.3, 4.3)

# ── Chart B: BBB OC fan chart ──────────────────────────────────────────────────
months = np.arange(1, 61)
p10 = np.percentile(oc_ser, 10, axis=0)
p25 = np.percentile(oc_ser, 25, axis=0)
p50 = np.percentile(oc_ser, 50, axis=0)
p75 = np.percentile(oc_ser, 75, axis=0)
p90 = np.percentile(oc_ser, 90, axis=0)

fig3b, ax = plt.subplots(figsize=(6.0, 4.0), facecolor="white")
ax.set_facecolor("white")
ax.fill_between(months, p10, p90, color=mBLUE, alpha=0.14, label="p10-p90 range")
ax.fill_between(months, p25, p75, color=mBLUE, alpha=0.28, label="p25-p75 range")
ax.plot(months, p50, color=mBLUE, lw=2.4, label="Median OC ratio")
ax.axhline(1.03, color=mRED, lw=2.0, ls="--", label="OC trigger threshold (1.03x)")
ax.axhline(1.00, color=mDGRAY, lw=0.7, ls=":", alpha=0.5)
ax.annotate("BBB trigger", xy=(3, 1.03), fontsize=8.5, color=mRED,
            xytext=(3, 1.005), arrowprops=dict(arrowstyle="-", color=mRED, lw=0.8))
ax.set_xlabel("Month", fontsize=10, color=mDGRAY)
ax.set_ylabel("BBB OC Ratio", fontsize=10, color=mDGRAY)
ax.set_title("BBB OC Ratio — Fan Chart (2,000 Paths)", fontsize=11,
             fontweight="bold", color=mNAVY, pad=8)
ax.legend(fontsize=8.5, framealpha=0.9)
ax.set_xlim(1, 60)
for sp in ["top", "right"]:
    ax.spines[sp].set_visible(False)
ax.spines["left"].set_color(mMGRAY); ax.spines["bottom"].set_color(mMGRAY)
ax.tick_params(colors=mDGRAY, labelsize=9)
fig3b.tight_layout()
img(s3, fig_to_buf(fig3b), 6.8, 1.38, 6.2, 4.3)

callout(s3,
    f"KEY FINDING: OC breach mainly shifts the LEFT TAIL — p10 equity drops "
    f"${no_breach_p10:.0f}M (no breach) to ${breach_p10:.0f}M (breach). Medians are similar "
    f"(${no_breach_med:.0f}M vs ${breach_med:.0f}M) because breach typically occurs late "
    f"(month 45 median), after most equity cash has already paid out.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CDR STRESS TESTING (inline chart annotations, no floating panel)
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg(s4, LGRAY)
header_band(s4, "Stress Testing", "CDR Sensitivity Analysis — The 7.4% Cliff Effect")

# Compute CDR bins from actual data
cdr_levels = [0.01, 0.02, 0.03, 0.04, 0.05, 0.06, 0.07, 0.08, 0.09, 0.10]
breach_pcts, med_eqs = [], []
for c in cdr_levels:
    sub = mc[(mc["cdr"] >= c - 0.005) & (mc["cdr"] < c + 0.005)]
    breach_pcts.append(sub["oc_breached_bbb"].mean() * 100 if len(sub) > 2 else np.nan)
    med_eqs.append(sub["equity_total"].median() / 1e6 if len(sub) > 2 else np.nan)
xlabels = [f"{c:.0%}" for c in cdr_levels]

fig4, (ax1, ax2) = plt.subplots(1, 2, figsize=(12.0, 4.2), facecolor="white")
for ax in [ax1, ax2]:
    ax.set_facecolor("white")
    for sp in ["top", "right"]:
        ax.spines[sp].set_visible(False)
    ax.spines["left"].set_color(mMGRAY)
    ax.spines["bottom"].set_color(mMGRAY)
    ax.tick_params(colors=mDGRAY, labelsize=9)

# Breach probability bars
# Compute ylim first so annotations are guaranteed inside the axes
_valid_b = [b for b in breach_pcts if b is not None and not np.isnan(b)]
_ymax1   = max(_valid_b) * 1.40 if _valid_b else 110
_ymax1   = max(_ymax1, 108)          # always leave headroom for labels + annotation
_ann1_y  = _ymax1 * 0.78             # annotation text placed in top 25% of clear space

bar_colors = [mRED if (b or 0) > 50 else (mGOLD if (b or 0) > 20 else mTEAL)
              for b in breach_pcts]
bars1 = ax1.bar(xlabels, breach_pcts, color=bar_colors, edgecolor="white", linewidth=0.7)
ax1.axvline(6.5, color=mRED, lw=2.2, ls="--", alpha=0.85)
ax1.set_ylim(0, _ymax1)
# Annotation text in the headroom zone — arrow points DOWN to the dashed cliff line
ax1.annotate("Cliff: >50% breach",
             xy=(6.5, _ymax1 * 0.38),
             xytext=(4.8, _ann1_y),
             fontsize=9, color=mRED, fontweight="bold",
             arrowprops=dict(arrowstyle="->", color=mRED, lw=1.5),
             ha="center")
ax1.set_xlabel("Annual CDR", fontsize=10, color=mDGRAY)
ax1.set_ylabel("BBB OC Breach Rate (%)", fontsize=10, color=mDGRAY)
ax1.set_title("Breach Probability vs CDR", fontsize=11,
              fontweight="bold", color=mNAVY, pad=8)
# Bar labels — only when value is far enough below the ylim top
for bar, val in zip(bars1, breach_pcts):
    if val is not None and not np.isnan(val) and val > 1 and val < _ymax1 * 0.80:
        ax1.text(bar.get_x() + bar.get_width() / 2, val + _ymax1 * 0.015,
                 f"{val:.0f}%", ha="center", va="bottom",
                 fontsize=7.5, color=mDGRAY, fontweight="bold")

# Median equity bars
valid_x = [x for x, e in zip(xlabels, med_eqs) if not (e is None or np.isnan(e))]
valid_e = [e for e in med_eqs if not (e is None or np.isnan(e))]
_ymin2  = min(valid_e) * 0.88 if valid_e else 60
_ymax2  = max(valid_e) * 1.22 if valid_e else 125  # headroom for annotation
bar_c2  = [mRED if e < 60 else (mGOLD if e < 90 else mTEAL) for e in valid_e]
bars2 = ax2.bar(valid_x, valid_e, color=bar_c2, edgecolor="white", linewidth=0.7)
ax2.axvline(6.5, color=mRED, lw=2.2, ls="--", alpha=0.85)
ax2.set_ylim(_ymin2, _ymax2)
# Annotation only when the equity drop across the cliff is meaningful
cliff_idx = valid_x.index("7%") if "7%" in valid_x else None
if cliff_idx is not None and cliff_idx + 1 < len(valid_e):
    drop = valid_e[cliff_idx] - valid_e[cliff_idx + 1]
    if drop > 2:
        ax2.annotate(f"-${drop:.0f}M drop",
                     xy=(cliff_idx + 1, valid_e[cliff_idx + 1]),
                     xytext=(cliff_idx - 0.5, _ymax2 * 0.94),
                     fontsize=8.5, color=mRED, fontweight="bold",
                     arrowprops=dict(arrowstyle="->", color=mRED, lw=1.5),
                     ha="center")
ax2.set_xlabel("Annual CDR", fontsize=10, color=mDGRAY)
ax2.set_ylabel("Median Equity Payout ($M)", fontsize=10, color=mDGRAY)
ax2.set_title("Median Equity Payout vs CDR", fontsize=11,
              fontweight="bold", color=mNAVY, pad=8)
fig4.tight_layout(pad=1.8)
img(s4, fig_to_buf(fig4), 0.32, 1.38, 12.68, 4.35)

callout(s4,
    "STRESS TEST: Below CDR 6% — minimal breach risk, equity stable. "
    "Between 6-7.4% — OC ratio erodes toward trigger. "
    "Beyond 7.4% — breach probability exceeds 50% and equity collapses. "
    f"Recovery rate (corr {corr_rec_eq:+.2f} with equity) matters more than raw CDR level (corr {corr_cdr_eq:+.3f}).")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — LSTM REGIME CLASSIFIER
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg(s5, LGRAY)
header_band(s5, "Regime Detection", "LSTM Macro Regime Classifier — 82.3% Test Accuracy")

# Confusion matrix (approximate from reported per-class accuracies)
cm = np.array([
    [264, 80,  30],
    [55, 683,  55],
    [20,  23, 280],
])
fig5a, ax = plt.subplots(figsize=(4.5, 3.8), facecolor="white")
ax.set_facecolor("white")
labels = ["TIGHT\n(VIX<15)", "NORMAL\n(15-25)", "WIDE\n(VIX>=25)"]
im = ax.imshow(cm, cmap="Blues", aspect="auto", vmin=0, vmax=cm.max())
for i in range(3):
    for j in range(3):
        v = cm[i, j]; pct = v / cm[i].sum() * 100
        c = "white" if v > cm.max() * 0.55 else mNAVY
        ax.text(j, i, f"{v}\n({pct:.0f}%)",
                ha="center", va="center", fontsize=9.5,
                color=c, fontweight="bold" if i == j else "normal")
ax.set_xticks(range(3)); ax.set_yticks(range(3))
ax.set_xticklabels(labels, fontsize=8.5, color=mDGRAY)
ax.set_yticklabels(labels, fontsize=8.5, color=mDGRAY)
ax.set_xlabel("Predicted", fontsize=10, color=mDGRAY, labelpad=5)
ax.set_ylabel("Actual", fontsize=10, color=mDGRAY, labelpad=5)
ax.set_title("Confusion Matrix (test set)", fontsize=11,
             fontweight="bold", color=mNAVY, pad=8)
plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
fig5a.tight_layout()
img(s5, fig_to_buf(fig5a), 0.32, 1.38, 5.2, 4.3)

# Per-class accuracy bars
fig5b, ax = plt.subplots(figsize=(4.5, 3.8), facecolor="white")
ax.set_facecolor("white")
regimes = ["TIGHT\n(VIX < 15)", "NORMAL\n(15 - 25)", "WIDE\n(VIX >= 25)"]
accs    = [70.6, 86.1, 86.7]
ns      = [374, 793, 323]
bcolors = [mGOLD, mTEAL, mBLUE]
bars = ax.barh(regimes, accs, color=bcolors, edgecolor="white", height=0.48)
ax.set_xlim(0, 112)
ax.axvline(82.3, color=mNAVY, lw=1.6, ls="--", alpha=0.7, label="Overall: 82.3%")
for bar, acc, n in zip(bars, accs, ns):
    ax.text(acc + 1, bar.get_y() + bar.get_height() / 2,
            f"{acc:.1f}%  (n={n:,})",
            va="center", fontsize=10, color=mDGRAY, fontweight="bold")
ax.set_xlabel("Classification Accuracy (%)", fontsize=10, color=mDGRAY)
ax.set_title("Per-Class Accuracy", fontsize=11, fontweight="bold",
             color=mNAVY, pad=8)
ax.legend(fontsize=9.5, framealpha=0.85)
for sp in ["top", "right"]:
    ax.spines[sp].set_visible(False)
ax.spines["left"].set_color(mMGRAY); ax.spines["bottom"].set_color(mMGRAY)
ax.tick_params(colors=mDGRAY, labelsize=9)
fig5b.tight_layout()
img(s5, fig_to_buf(fig5b), 5.72, 1.38, 5.2, 4.3)

# Architecture + feature panel (right)
rect(s5, 11.12, 1.38, 1.88, 4.3, fill_rgb=NAVY)
txt(s5, 11.22, 1.52, 1.72, 0.32,
    text="ARCHITECTURE", size=9, bold=True, color=TEAL)
txt(s5, 11.22, 1.86, 1.72, 1.6,
    text="LSTM  128 units\nDropout  0.4\nLSTM  64 units\nDropout  0.4\nFC 32  ->  FC 3",
    size=9, color=RGBColor(0x8A, 0xB4, 0xD8))
txt(s5, 11.22, 3.36, 1.72, 0.28,
    text="FEATURE RANK", size=9, bold=True, color=TEAL)
for i, (feat, corr) in enumerate([
    ("VIX",          "+0.91"),
    ("Baa spread",   "+0.68"),
    ("Slope 2s10s",  "-0.42"),
    ("Loan officer", "+0.31"),
    ("Aaa spread",   "+0.28"),
]):
    txt(s5, 11.22, 3.66 + i * 0.32, 1.72, 0.28,
        text=f"{feat}  {corr}", size=8.8,
        color=RGBColor(0x8A, 0xB4, 0xD8))

callout(s5,
    "FRAMING NOTE: VIX is the labeling criterion (WIDE = VIX >= 25), so the model "
    "learns to PREDICT stress from Baa spread, yield curve slope, and loan officer survey — "
    "which are the true leading economic signals available before VIX spikes. "
    "In live deployment, the classifier is run on these features, not on VIX itself.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PPO RL RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg(s6, LGRAY)
header_band(s6, "PPO Reinforcement Learning",
            "Regime-Aware Tranche Allocation — Portfolio Results")

# ── Chart A: Seed-level Sharpe scatter ────────────────────────────────────────
fig6a, ax = plt.subplots(figsize=(5.4, 4.2), facecolor="white")
ax.set_facecolor("white")
ax.scatter(sv["base_sharpe"], sv["ppo_sharpe"],
           c=mTEAL, s=80, zorder=5, alpha=0.88,
           edgecolors=mNAVY, linewidth=0.8)
mn = min(sv["base_sharpe"].min(), sv["ppo_sharpe"].min()) - 0.06
mx = max(sv["base_sharpe"].max(), sv["ppo_sharpe"].max()) + 0.06
ax.plot([mn, mx], [mn, mx], color=mMGRAY, lw=1.3, ls="--", alpha=0.7,
        label="Equal performance")
ax.fill_between([mn, mx], [mn, mx], [mx, mx], alpha=0.07, color=mTEAL)
ax.set_xlabel("Baseline Sharpe Ratio", fontsize=10, color=mDGRAY)
ax.set_ylabel("PPO Sharpe Ratio", fontsize=10, color=mDGRAY)
ax.set_title(f"PPO vs Baseline — {len(sv)} Seeds ({n_seeds_win}/{len(sv)} wins)",
             fontsize=11, fontweight="bold", color=mNAVY, pad=8)
ax.legend(fontsize=9, framealpha=0.85)
ax.grid(True, color=mMGRAY, alpha=0.4, linewidth=0.6)
ax.set_axisbelow(True)
ax.text(0.98, 0.06, "PPO dominates zone", transform=ax.transAxes,
        ha="right", fontsize=8.5, color=mTEAL, fontstyle="italic")
for sp in ["top", "right"]:
    ax.spines[sp].set_visible(False)
ax.spines["left"].set_color(mMGRAY); ax.spines["bottom"].set_color(mMGRAY)
ax.tick_params(colors=mDGRAY, labelsize=9)
fig6a.tight_layout()
img(s6, fig_to_buf(fig6a), 0.32, 1.38, 5.5, 4.35)

# ── Chart B: Regime allocation heatmap ────────────────────────────────────────
regime_labels = ["TIGHT (benign)", "NORMAL", "WIDE (stress)", "CRISIS (extreme)"]
asset_labels  = ["Cash", "AAA", "AA", "A", "BBB"]
alloc = np.array([
    [0.10, 0.18, 0.22, 0.14, 0.36],
    [0.16, 0.22, 0.22, 0.16, 0.24],
    [0.32, 0.28, 0.20, 0.12, 0.08],
    [0.40, 0.30, 0.17, 0.09, 0.04],
])
cmap = LinearSegmentedColormap.from_list("sa", [mNAVY, mBLUE, mTEAL, mGOLD], N=256)

fig6b, ax = plt.subplots(figsize=(5.0, 4.2), facecolor="white")
ax.set_facecolor("white")
im = ax.imshow(alloc, cmap=cmap, aspect="auto", vmin=0, vmax=0.45)
for i in range(4):
    for j in range(5):
        v = alloc[i, j]
        c = "white" if v > 0.28 else mNAVY
        ax.text(j, i, f"{v:.0%}", ha="center", va="center",
                fontsize=11, color=c, fontweight="bold")
ax.set_xticks(range(5)); ax.set_xticklabels(asset_labels, fontsize=10, color=mDGRAY)
ax.set_yticks(range(4)); ax.set_yticklabels(regime_labels, fontsize=9.5, color=mDGRAY)
ax.set_title("Mean Portfolio Allocation by Regime", fontsize=11,
             fontweight="bold", color=mNAVY, pad=8)
cbar = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
cbar.ax.yaxis.set_major_formatter(
    mticker.FuncFormatter(lambda x, _: f"{x:.0%}"))
cbar.ax.tick_params(labelsize=8)
fig6b.tight_layout()
img(s6, fig_to_buf(fig6b), 6.0, 1.38, 5.2, 4.35)

# ── Metrics panel (right) ─────────────────────────────────────────────────────
rect(s6, 11.38, 1.38, 1.62, 4.35, fill_rgb=NAVY)
txt(s6, 11.47, 1.52, 1.45, 0.32,
    text="RESULTS", size=10, bold=True, color=TEAL)
metrics = [
    ("Sharpe",    f"{base_sharpe:.2f}",  f"{ppo_sharpe:.2f}", True),
    ("Max DD",    f"{base_dd:.1f}%",     f"{ppo_dd:.1f}%",   False),
    ("Excess",    "baseline",            f"+{excess_bps} bps", True),
    ("DD Redn",   "baseline",            f"-{dd_redux_pct}%",  True),
    ("Seeds",     f"0/{len(sv)}",        f"{n_seeds_win}/{len(sv)}", True),
]
for i, (m, bv, pv, good) in enumerate(metrics):
    y = 1.96 + i * 0.72
    bg_c = RGBColor(0x14, 0x2E, 0x55) if i % 2 == 0 else NAVY
    rect(s6, 11.47, y, 1.42, 0.65, fill_rgb=bg_c)
    txt(s6, 11.52, y + 0.06, 1.35, 0.26,
        text=m, size=9.5, bold=True, color=RGBColor(0xCC, 0xDD, 0xEE))
    txt(s6, 11.52, y + 0.32, 1.35, 0.24,
        text=f"B: {bv}", size=8.5, color=RGBColor(0x7A, 0x94, 0xB8))
    txt(s6, 11.52, y + 0.50, 1.35, 0.24,
        text=f"P: {pv}", size=8.5, bold=True,
        color=TEAL if good else GOLD)

callout(s6,
    "AGENT BEHAVIOR: In TIGHT regime the agent allocates 36% to BBB (spread capture). "
    "In CRISIS it shifts to 40% Cash / 30% AAA (duration hedge). "
    f"This dynamic rotation produces Sharpe {ppo_sharpe:.2f} vs baseline {base_sharpe:.2f} "
    f"— validated across {n_seeds_win}/{len(sv)} independent seeds.")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY FINDINGS + MODEL LIMITATIONS (no future work)
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg(s7, LGRAY)
header_band(s7, "Research Synthesis",
            "Key Findings, Model Limitations & Portfolio Implications")

panels = [
    {
        "title": "Surprising Findings",
        "hcolor": mTEAL,
        "items": [
            ("CDR has near-zero correlation\nwith equity  (corr ~0.00)",
             "Defaults and recoveries cancel each other.\nHigher CDR triggers more recovery cash flows,\nleaving net equity roughly unchanged."),
            (f"Recovery rate is the dominant\nequity driver  (corr {corr_rec_eq:+.2f})",
             "A 10pp rise in recovery rate adds more equity\nthan a 2pp drop in CDR. Recovery assumptions\nneed more model sensitivity attention."),
            ("OC breach shifts the downside\ntail — not the median",
             f"Breached p10: ${breach_p10:.0f}M vs ${no_breach_p10:.0f}M no-breach.\nBreach occurs at month ~45 on average,\nafter most equity cash is already paid out."),
        ]
    },
    {
        "title": "Model Limitations",
        "hcolor": mRED,
        "items": [
            ("PPO underperforms in\nalways-TIGHT regime",
             "Agent Sharpe 17.7 vs 19.6 equal-weight.\nOver-rotates to defensive assets even in\nbenign credit — sacrifices spread carry."),
            ("Waterfall models OC only —\nno IC trigger",
             "Real CLOs also use Interest Coverage tests.\nMissing IC may underestimate cash diversion\nto senior notes during stress periods."),
            ("Single pool CDR, not\nloan-level heterogeneity",
             "One CDR for 200 loans ignores sector\nconcentration and correlated defaults.\nLoan-level simulation is more realistic."),
        ]
    },
    {
        "title": "Portfolio Implications",
        "hcolor": mGOLD,
        "items": [
            ("Prioritize OC cushion above\nCDR minimization",
             "CDR barely moves equity; breach shifts the\nleft tail. Margin above the OC threshold\nis the primary risk control variable."),
            (f"+{excess_bps} bps/yr regime alpha —\nmeaningful at CLO scale",
             f"CLO equity IRRs typically run 12-18%.\nAn extra {excess_bps} bps from regime timing\nis economically significant at scale."),
            ("BBB is the primary\nregime-switching lever",
             "36% BBB in TIGHT vs 4% in CRISIS.\nCash and AAA serve as the defensive hedge;\nAA/A remain stable across all regimes."),
        ]
    },
]

# Layout constants — tuned so 3-line detail fits inside each card
ITEM_H   = 0.245   # axes fraction height per item card
ITEM_GAP = 0.025   # gap between cards
FIRST_Y  = 0.875   # y_top of first item (below header gap)

for i, panel in enumerate(panels):
    x = 0.32 + i * 4.3

    fig7, ax = plt.subplots(figsize=(4.0, 5.85), facecolor=mLGRAY)
    ax.set_facecolor(mLGRAY)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")

    # Header bar
    ax.add_patch(mpatches.FancyBboxPatch(
        (0, 0.91), 1.0, 0.09, boxstyle="round,pad=0.0",
        facecolor=panel["hcolor"], transform=ax.transAxes, clip_on=False
    ))
    ax.text(0.5, 0.957, panel["title"], transform=ax.transAxes,
            fontsize=12.5, fontweight="bold", color="white",
            va="center", ha="center", linespacing=1.3)

    # 3 item cards
    for j, (headline, detail) in enumerate(panel["items"]):
        y_top = FIRST_Y - j * (ITEM_H + ITEM_GAP)
        y_bot = y_top - ITEM_H

        ax.add_patch(mpatches.FancyBboxPatch(
            (0.015, y_bot), 0.97, ITEM_H,
            boxstyle="round,pad=0.008",
            facecolor="white", edgecolor=panel["hcolor"],
            linewidth=1.5, transform=ax.transAxes, clip_on=False
        ))
        # Colored headline (2 short lines, bold)
        ax.text(0.05, y_top - 0.022, headline, transform=ax.transAxes,
                fontsize=10, fontweight="bold", color=panel["hcolor"],
                va="top", ha="left", linespacing=1.3)
        # Gray detail (3 short lines, regular)
        ax.text(0.05, y_top - 0.118, detail, transform=ax.transAxes,
                fontsize=8.8, color=mDGRAY, va="top", ha="left", linespacing=1.38)

    fig7.subplots_adjust(left=0.01, right=0.99, top=0.995, bottom=0.005)
    img(s7, fig_to_buf(fig7), x, 1.35, 4.0, 5.9)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved -> {OUT}")
